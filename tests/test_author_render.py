"""Tests for papercast.author.render — assembling the lecture PPTX from
slides_plan.json + figures.json + lab_template.pptx.

Uses the real template + the test paper's hand-authored slides_plan as
the fixture. The actual visual quality is something only a human can
judge; these tests lock the structural contract (right layouts, right
placeholders filled, all referenced images present).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from papercast.author.render import (
    PageSpec,
    SlidesPlan,
    assemble_pptx,
    load_slides_plan,
    parse_script_md,
)

REPO = Path(__file__).resolve().parents[1]
TEMPLATE = REPO / "templates" / "lab_template.pptx"
FIXTURE_PLAN = REPO / "work" / "e8f6731a14" / "slides_plan.json"
FIXTURE_FIGURES_DIR = REPO / "work" / "e8f6731a14" / "figures"


@pytest.fixture(scope="module")
def assembled(tmp_path_factory: pytest.TempPathFactory) -> tuple[Path, SlidesPlan]:
    if not FIXTURE_PLAN.exists():
        pytest.skip("slides_plan.json missing; produce the fixture first")
    if not TEMPLATE.exists():
        pytest.skip("lab_template.pptx missing")
    plan = load_slides_plan(FIXTURE_PLAN)
    out = tmp_path_factory.mktemp("pptx") / "out.pptx"
    assemble_pptx(plan, TEMPLATE, FIXTURE_FIGURES_DIR, out)
    return out, plan


def test_load_slides_plan_basic() -> None:
    if not FIXTURE_PLAN.exists():
        pytest.skip("plan missing")
    plan = load_slides_plan(FIXTURE_PLAN)
    assert plan.paper_id == "e8f6731a14"
    assert plan.total_pages == len(plan.pages)
    assert plan.total_pages >= 10


def test_assembled_pptx_exists_and_has_correct_slide_count(
    assembled: tuple[Path, SlidesPlan]
) -> None:
    out_path, plan = assembled
    assert out_path.exists()
    prs = Presentation(out_path)
    assert len(prs.slides) == plan.total_pages


def test_each_slide_uses_the_planned_layout(
    assembled: tuple[Path, SlidesPlan]
) -> None:
    out_path, plan = assembled
    prs = Presentation(out_path)
    for slide, page in zip(prs.slides, plan.pages, strict=True):
        assert slide.slide_layout.name == page.layout, (
            f"page {page.page_no} expected layout {page.layout!r} "
            f"but got {slide.slide_layout.name!r}"
        )


def test_text_placeholders_get_filled(assembled: tuple[Path, SlidesPlan]) -> None:
    """Every non-image field in the plan must show up as text on the
    corresponding slide's placeholder."""
    out_path, plan = assembled
    prs = Presentation(out_path)
    for slide, page in zip(prs.slides, plan.pages, strict=True):
        # Map placeholder name -> rendered text on the slide
        slide_phs: dict[str, str] = {}
        for shp in slide.placeholders:
            if shp.has_text_frame:
                slide_phs[shp.name] = shp.text_frame.text

        for field_name, field_value in page.fields.items():
            if field_name == "Image":
                continue
            if field_name not in slide_phs:
                continue  # Layout doesn't have this field — caught elsewhere
            rendered = slide_phs[field_name]
            if isinstance(field_value, list):
                # Each bullet should appear in the rendered text.
                for bullet in field_value:
                    assert bullet in rendered, (
                        f"page {page.page_no} {field_name}: "
                        f"bullet {bullet!r} missing from {rendered!r}"
                    )
            else:
                assert str(field_value) in rendered, (
                    f"page {page.page_no} {field_name}: "
                    f"value {field_value!r} missing from {rendered!r}"
                )


def _slide_has_image(slide) -> bool:
    """A slide carries an image either as a free Picture shape (shape_type
    == PICTURE / 13) or as a PICTURE placeholder that has been filled
    (shape_type stays PLACEHOLDER / 14, but `.image` is present)."""
    for shp in slide.shapes:
        if shp.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            return True
        if shp.is_placeholder:
            try:
                if shp.image is not None:
                    return True
            except Exception:  # noqa: BLE001 — placeholder w/o image raises
                continue
    return False


def test_image_placeholders_get_pictures(
    assembled: tuple[Path, SlidesPlan]
) -> None:
    """Every page that specifies an Image must have actual image content
    on the resulting slide (either a Picture shape or a filled PICTURE
    placeholder)."""
    out_path, plan = assembled
    prs = Presentation(out_path)
    for slide, page in zip(prs.slides, plan.pages, strict=True):
        if "Image" not in page.fields:
            continue
        assert _slide_has_image(slide), (
            f"page {page.page_no} missing image for "
            f"Image={page.fields['Image']!r}"
        )


def test_bullets_render_as_multiple_paragraphs(
    assembled: tuple[Path, SlidesPlan]
) -> None:
    """A list value should produce one paragraph per bullet, not a
    single concatenated line."""
    out_path, plan = assembled
    prs = Presentation(out_path)
    for slide, page in zip(prs.slides, plan.pages, strict=True):
        for field_name, field_value in page.fields.items():
            if not isinstance(field_value, list):
                continue
            ph = next((shp for shp in slide.placeholders
                       if shp.name == field_name), None)
            if ph is None or not ph.has_text_frame:
                continue
            paragraphs = [p.text for p in ph.text_frame.paragraphs]
            assert len(paragraphs) == len(field_value), (
                f"page {page.page_no} {field_name}: expected "
                f"{len(field_value)} paragraphs, got {len(paragraphs)}"
            )


def test_unknown_layout_fails_loudly(tmp_path: Path) -> None:
    """If the plan references a layout the template doesn't have, we
    fail fast — silent fallback would corrupt the lab format."""
    if not TEMPLATE.exists():
        pytest.skip("template missing")
    bad_plan = SlidesPlan(
        paper_id="x",
        total_pages=1,
        target_duration_sec=60,
        pages=[PageSpec(page_no=1, layout="DoesNotExist", fields={"Title": "x"})],
    )
    out = tmp_path / "x.pptx"
    with pytest.raises(ValueError, match="layout.*not found"):
        assemble_pptx(bad_plan, TEMPLATE, tmp_path, out)


def test_missing_image_fails_loudly(tmp_path: Path) -> None:
    """If the plan references an Image id that figures.json / disk
    doesn't have, we fail fast."""
    if not TEMPLATE.exists():
        pytest.skip("template missing")
    fig_dir = tmp_path / "figs"
    fig_dir.mkdir()
    (fig_dir / "figures.json").write_text("[]", encoding="utf-8")
    plan = SlidesPlan(
        paper_id="x",
        total_pages=1,
        target_duration_sec=60,
        pages=[PageSpec(
            page_no=1,
            layout="JournalIntro",
            fields={"Subtitle": "X", "Bullets": ["a"], "Image": "fig_999"},
        )],
    )
    out = tmp_path / "x.pptx"
    with pytest.raises(ValueError, match="image.*not found|fig_999"):
        assemble_pptx(plan, TEMPLATE, fig_dir, out)


def test_real_pptx_smoke_run(assembled: tuple[Path, SlidesPlan]) -> None:
    """Full end-to-end: opening the assembled pptx round-trips without
    errors and slide count matches plan. Also, file size is plausible."""
    out_path, plan = assembled
    size = out_path.stat().st_size
    # Template is ~534 KB; with images embedded we expect ~600 KB or more.
    assert size > 200_000, f"assembled pptx suspiciously small ({size} bytes)"
    Presentation(out_path)  # round-trip


# ---------------------------------------------------------------------------
# script.md → speaker notes
# ---------------------------------------------------------------------------


SAMPLE_SCRIPT_MD = """\
# Title (header is ignored)

> 总页数: 3

## Page 1
本次报告的主题是 X。

## Page 2
首先看背景。
然后看方法。

## Page 3
谢谢。
"""


def test_parse_script_md_extracts_per_page(tmp_path: Path) -> None:
    p = tmp_path / "script.md"
    p.write_text(SAMPLE_SCRIPT_MD, encoding="utf-8")
    notes = parse_script_md(p)
    assert notes == {
        1: "本次报告的主题是 X。",
        2: "首先看背景。\n然后看方法。",
        3: "谢谢。",
    }


def test_parse_script_md_missing_file_returns_empty(tmp_path: Path) -> None:
    notes = parse_script_md(tmp_path / "does_not_exist.md")
    assert notes == {}


def test_parse_script_md_strips_trailing_metadata_fence(tmp_path: Path) -> None:
    """Regression: the LLM Scripter appends a `---\\ntotal_chars: ...`
    fence at the end of script.md. Without filtering, this metadata
    bleeds into the LAST page's speaker notes (and gets read aloud by
    TTS). The parser must drop everything from the last `---` line on
    the last page."""
    content = (
        "## Page 1\n"
        "封面页讲稿。\n"
        "\n"
        "## Page 2\n"
        "结束页讲稿。\n"
        "\n"
        "---\n"
        "total_chars: 12\n"
        "estimated_seconds: 4\n"
        "in_target_range: false\n"
    )
    p = tmp_path / "script.md"
    p.write_text(content, encoding="utf-8")
    notes = parse_script_md(p)
    assert notes == {
        1: "封面页讲稿。",
        2: "结束页讲稿。",
    }


def test_parse_script_md_keeps_internal_horizontal_rules(tmp_path: Path) -> None:
    """A `---` that appears INSIDE a page's body (not at the trailing
    metadata fence) should be preserved — it might be a real horizontal
    rule the speaker wants in their notes."""
    content = (
        "## Page 1\n"
        "第一段。\n"
        "---\n"
        "第二段。\n"
        "\n"
        "## Page 2\n"
        "结尾。\n"
    )
    p = tmp_path / "script.md"
    p.write_text(content, encoding="utf-8")
    notes = parse_script_md(p)
    # Page 1 keeps its internal `---`; page 2 has no fence so nothing stripped.
    assert "---" in notes[1]
    assert notes[2] == "结尾。"


# ---------------------------------------------------------------------------
# Bullets adaptive font size (Fix-C r1 + Fix-7 r2)
# ---------------------------------------------------------------------------


def _build_bullets_pptx(
    tmp_path: Path, layout_name: str, n_bullets: int,
) -> Path:
    """Assemble a 1-page deck with `n_bullets` items in the named layout."""
    plan = SlidesPlan(
        paper_id="bullets_test",
        total_pages=1,
        target_duration_sec=480,
        pages=[
            PageSpec(
                page_no=1,
                layout=layout_name,
                fields={
                    "Subtitle": "测试",
                    "Bullets": [f"项目 {i+1}" for i in range(n_bullets)],
                },
            ),
        ],
    )
    out = tmp_path / "out.pptx"
    figures_dir = tmp_path / "figures"
    figures_dir.mkdir()
    (figures_dir / "figures.json").write_text("[]", encoding="utf-8")

    from papercast.author.render import assemble_pptx
    assemble_pptx(plan, TEMPLATE, figures_dir, out)
    return out


def _read_bullets_size(pptx_path: Path, layout_name: str) -> set[float]:
    """Return the set of pt sizes seen across all bullets paragraphs."""
    pres = Presentation(pptx_path)
    layout = next(l for l in pres.slide_layouts if l.name == layout_name)
    bullets_idx = next(
        ph.placeholder_format.idx for ph in layout.placeholders if ph.name == "Bullets"
    )
    slide = list(pres.slides)[0]
    ph = next(p for p in slide.placeholders if p.placeholder_format.idx == bullets_idx)

    sizes: set[float] = set()
    for para in ph.text_frame.paragraphs:
        if not (para.runs or para.text):
            continue
        size = None
        for run in para.runs:
            if run.font.size:
                size = run.font.size
                break
        if size is None:
            size = para.font.size
        if size is not None:
            sizes.add(size.pt)
    return sizes


# Layout-aware sizing matrix (r2). Tier ceilings (heights are read via
# python-pptx's Length.cm — see _bullets_size_ceiling docstring):
#   Bullets H ≥ 13 cm  → 24 pt   (text-heavy: TextOnly / Discussion /
#                                 Results / Background / Methods)
#   9 ≤ H < 13         → 22 pt   (TOC: 11.20 cm)
#   5 ≤ H < 9          → 20 pt   (WideImage: 6.4-7.8 cm)
#   H < 5              → 18 pt   (JournalIntro: 3.40 cm)
# Each tier steps DOWN by 2 pt for every extra 2 paragraphs past 5.

@pytest.mark.parametrize("layout, n_bullets, expected_pt", [
    # Tier 1 — text-heavy container (Background_TextOnly = 15.20 cm)
    ("Background_TextOnly", 3, 24),
    ("Background_TextOnly", 5, 24),
    ("Background_TextOnly", 6, 22),
    ("Background_TextOnly", 8, 20),
    ("Background_TextOnly", 10, 18),

    # Tier 2 — medium container (TOC = 11.20 cm)
    ("TOC", 4, 22),
    ("TOC", 6, 20),

    # Tier 3 — image-and-text container (Methods_WideImage = 7.40 cm)
    ("Methods_WideImage", 4, 20),
    ("Methods_WideImage", 6, 18),
    ("Methods_WideImage", 9, 16),

    # Tier 4 — tight container (JournalIntro = 3.40 cm)
    ("JournalIntro", 3, 18),
    ("JournalIntro", 7, 16),
])
def test_bullets_font_size_layout_aware(
    tmp_path_factory: pytest.TempPathFactory,
    layout: str, n_bullets: int, expected_pt: int,
) -> None:
    """The font-size schedule must take the placeholder's height into
    account, not just the bullet count, so 5 bullets in a 6 cm-tall
    Bullets box (Background_TextOnly) read larger than 5 bullets in a
    3 cm-tall box (Methods_WideImage)."""
    if not TEMPLATE.exists():
        pytest.skip("lab_template.pptx missing")

    out = _build_bullets_pptx(
        tmp_path_factory.mktemp(f"{layout}_{n_bullets}"),
        layout_name=layout,
        n_bullets=n_bullets,
    )
    sizes = _read_bullets_size(out, layout)
    assert sizes == {float(expected_pt)}, (
        f"layout={layout} n={n_bullets}: expected {expected_pt}pt, got {sizes}"
    )


def test_bullets_size_floor_at_12pt(tmp_path_factory: pytest.TempPathFactory) -> None:
    """Even on a tight container with many bullets the floor is 12 pt.

    Note: with the r2 schedule this is hard to trigger naturally —
    JournalIntro tier is 18 pt, so n=20 gives 18−6=12 pt → exactly the
    floor. For a tighter test we'd need more bullets, but n=20 already
    overflows the placeholder visually so we stop there.
    """
    if not TEMPLATE.exists():
        pytest.skip("lab_template.pptx missing")
    out = _build_bullets_pptx(
        tmp_path_factory.mktemp("floor"),
        layout_name="JournalIntro",  # 3.40 cm tier (18 pt ceiling)
        n_bullets=20,                # base − 6 = 12 pt
    )
    sizes = _read_bullets_size(out, "JournalIntro")
    assert sizes == {12.0}


def test_bullets_size_text_only_layouts_use_largest_tier(
    tmp_path_factory: pytest.TempPathFactory,
) -> None:
    """Sanity check: every *_TextOnly + Discussion + Results layout in
    the lab template falls into the 24 pt tier so they all start at
    24 pt for ≤5 bullets."""
    if not TEMPLATE.exists():
        pytest.skip("lab_template.pptx missing")
    for layout_name in (
        "Background_TextOnly",
        "Methods_TextOnly",
        "Experiment_TextOnly",
        "Discussion",
        "Results",
    ):
        out = _build_bullets_pptx(
            tmp_path_factory.mktemp(f"large_{layout_name}"),
            layout_name=layout_name, n_bullets=4,
        )
        sizes = _read_bullets_size(out, layout_name)
        assert sizes == {24.0}, f"{layout_name}: expected 24pt, got {sizes}"


def test_bullets_autosize_fallback_enabled(
    tmp_path_factory: pytest.TempPathFactory,
) -> None:
    """Even with an explicit font size, the Bullets text frame must have
    auto_size enabled as a safety net for unusually long single-bullet
    text."""
    if not TEMPLATE.exists():
        pytest.skip("lab_template.pptx missing")
    from pptx.enum.text import MSO_AUTO_SIZE

    plan = SlidesPlan(
        paper_id="autosize_test",
        total_pages=1,
        target_duration_sec=480,
        pages=[
            PageSpec(
                page_no=1,
                layout="Background_TextOnly",
                fields={
                    "Subtitle": "测试",
                    "Bullets": ["一条很长很长很长很长很长很长很长很长很长的描述" * 3],
                },
            ),
        ],
    )
    out = tmp_path_factory.mktemp("autosize") / "out.pptx"
    figures_dir = out.parent / "figures"
    figures_dir.mkdir()
    (figures_dir / "figures.json").write_text("[]", encoding="utf-8")

    from papercast.author.render import assemble_pptx
    assemble_pptx(plan, TEMPLATE, figures_dir, out)

    pres = Presentation(out)
    layout = next(l for l in pres.slide_layouts if l.name == "Background_TextOnly")
    bullets_idx = next(
        ph.placeholder_format.idx for ph in layout.placeholders if ph.name == "Bullets"
    )
    bullets_ph = next(
        ph for ph in list(pres.slides)[0].placeholders
        if ph.placeholder_format.idx == bullets_idx
    )
    assert bullets_ph.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    assert bullets_ph.text_frame.word_wrap is True


# Pure unit test for the height → ceiling mapping (no PPTX involved).
def test_bullets_size_ceiling_tiers() -> None:
    from papercast.author.render import _bullets_size_ceiling

    # Tier 1: ≥ 13 cm (text-heavy layouts in the lab template)
    assert _bullets_size_ceiling(13.0) == 24
    assert _bullets_size_ceiling(15.20) == 24      # Background_TextOnly
    assert _bullets_size_ceiling(15.93) == 24      # Discussion / Results

    # Tier 2: 9 - 13 cm (TOC)
    assert _bullets_size_ceiling(9.0) == 22
    assert _bullets_size_ceiling(11.20) == 22      # TOC
    assert _bullets_size_ceiling(12.99) == 22

    # Tier 3: 5 - 9 cm (image-and-text layouts)
    assert _bullets_size_ceiling(5.0) == 20
    assert _bullets_size_ceiling(6.40) == 20       # Background_WideImage
    assert _bullets_size_ceiling(7.40) == 20       # Methods_WideImage
    assert _bullets_size_ceiling(7.80) == 20       # Experiment_WideImage

    # Tier 4: < 5 cm (JournalIntro)
    assert _bullets_size_ceiling(3.40) == 18       # JournalIntro
    assert _bullets_size_ceiling(0.0) == 18


def test_assembled_slide_notes_match_script(
    tmp_path_factory: pytest.TempPathFactory,
) -> None:
    """When page_notes is supplied, each slide's speaker-notes pane should
    contain the corresponding page's script."""
    if not FIXTURE_PLAN.exists() or not TEMPLATE.exists():
        pytest.skip("fixtures missing")
    plan = load_slides_plan(FIXTURE_PLAN)
    notes = {p.page_no: f"NOTES FOR PAGE {p.page_no}" for p in plan.pages}
    out = tmp_path_factory.mktemp("pptx_notes") / "out.pptx"
    assemble_pptx(plan, TEMPLATE, FIXTURE_FIGURES_DIR, out, page_notes=notes)
    prs = Presentation(out)
    for slide, page in zip(prs.slides, plan.pages, strict=True):
        rendered = slide.notes_slide.notes_text_frame.text
        assert f"NOTES FOR PAGE {page.page_no}" in rendered


def test_assembled_without_notes_has_empty_notes(
    assembled: tuple[Path, SlidesPlan]
) -> None:
    """When page_notes is omitted, slides have no speaker notes."""
    out_path, plan = assembled
    prs = Presentation(out_path)
    for slide in prs.slides:
        # Notes slide may or may not exist; if it does, should be empty.
        if slide.has_notes_slide:
            assert slide.notes_slide.notes_text_frame.text.strip() == ""


# ---------------------------------------------------------------------------
# template_vars substitution (e.g. {{REPORT_DATE}} → "2026-05-29")
# ---------------------------------------------------------------------------


def test_template_vars_substitute_in_text_fields(
    tmp_path_factory: pytest.TempPathFactory,
) -> None:
    """A plan field value of "{{REPORT_DATE}}" must be replaced with the
    value supplied in template_vars before being written to the slide."""
    if not FIXTURE_PLAN.exists() or not TEMPLATE.exists():
        pytest.skip("fixtures missing")
    plan = load_slides_plan(FIXTURE_PLAN)
    out = tmp_path_factory.mktemp("vars") / "out.pptx"
    assemble_pptx(
        plan, TEMPLATE, FIXTURE_FIGURES_DIR, out,
        template_vars={"REPORT_DATE": "2026-05-29"},
    )
    prs = Presentation(out)
    cover = prs.slides[0]
    cover_text = "\n".join(
        shp.text_frame.text for shp in cover.placeholders if shp.has_text_frame
    )
    assert "2026-05-29" in cover_text
    assert "{{REPORT_DATE}}" not in cover_text


def test_template_vars_unset_leaves_placeholder(
    tmp_path_factory: pytest.TempPathFactory,
) -> None:
    """When template_vars is omitted, the literal "{{REPORT_DATE}}" should
    survive — that's the signal that the human review hasn't filled the
    date yet."""
    if not FIXTURE_PLAN.exists() or not TEMPLATE.exists():
        pytest.skip("fixtures missing")
    plan = load_slides_plan(FIXTURE_PLAN)
    out = tmp_path_factory.mktemp("vars") / "out.pptx"
    assemble_pptx(plan, TEMPLATE, FIXTURE_FIGURES_DIR, out)
    prs = Presentation(out)
    cover_text = "\n".join(
        shp.text_frame.text for shp in prs.slides[0].placeholders
        if shp.has_text_frame
    )
    assert "{{REPORT_DATE}}" in cover_text


def test_template_vars_substitute_inside_lists(tmp_path: Path) -> None:
    """Substitution should also work inside list values (Bullets)."""
    if not TEMPLATE.exists():
        pytest.skip("template missing")
    plan = SlidesPlan(
        paper_id="x", total_pages=1, target_duration_sec=60,
        pages=[PageSpec(
            page_no=1, layout="TOC",
            fields={"Title": "Title for {{NAME}}", "Bullets": ["Hello {{NAME}}!"]},
        )],
    )
    out = tmp_path / "x.pptx"
    figs = tmp_path / "figs"
    figs.mkdir()
    (figs / "figures.json").write_text("[]", encoding="utf-8")
    assemble_pptx(plan, TEMPLATE, figs, out, template_vars={"NAME": "alice"})
    prs = Presentation(out)
    text = "\n".join(
        shp.text_frame.text for shp in prs.slides[0].placeholders
        if shp.has_text_frame
    )
    assert "Title for alice" in text
    assert "Hello alice!" in text
