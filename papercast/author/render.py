"""Assemble the lecture .pptx from a slides_plan.json + figures + template.

This is the JSON-first assembler called for in §6.3.2 of the design doc:
  - LLM (or hand-author) produces slides_plan.json with one entry per page,
    naming the layout and providing field values keyed by placeholder name.
  - This module opens the lab template, copies one slide per plan page using
    the named layout, and writes each field into its placeholder.
  - No layout decisions are made here — the template is the schema, this is
    pure plumbing.

Lookup keys throughout are placeholder NAMES (not idx), because that's what
slides_plan.json uses and what the meta.json contract exposes.

Visual conventions (also live as code rather than master decisions until
the lab template is updated to match):
  - Cover Title and Title_chinese are centered.
  - Subtitle (small label upper-left of every content layout) is bold.
  - Image fields use *contain* fit — keep aspect, never crop or stretch.
"""

from __future__ import annotations

import json
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.slide import SlideLayout
from pptx.util import Pt


@dataclass(frozen=True)
class PageSpec:
    page_no: int
    layout: str
    fields: dict[str, Any] = field(default_factory=dict)


@dataclass
class SlidesPlan:
    paper_id: str
    total_pages: int
    target_duration_sec: int
    pages: list[PageSpec] = field(default_factory=list)


def load_slides_plan(path: Path) -> SlidesPlan:
    payload = json.loads(Path(path).read_text(encoding="utf-8"))
    pages = [
        PageSpec(
            page_no=p["page_no"],
            layout=p["layout"],
            fields=dict(p.get("fields", {})),
        )
        for p in payload["pages"]
    ]
    return SlidesPlan(
        paper_id=payload["paper_id"],
        total_pages=payload["total_pages"],
        target_duration_sec=payload.get("target_duration_sec", 480),
        pages=pages,
    )


def assemble_pptx(
    plan: SlidesPlan,
    template_path: Path,
    figures_dir: Path,
    out_path: Path,
    page_notes: dict[int, str] | None = None,
    template_vars: dict[str, str] | None = None,
) -> None:
    """Build the .pptx from `plan`, using `template_path` as the schema and
    `figures_dir` (containing figures.json + PNGs) as the image source.

    `page_notes` (optional) maps page_no -> spoken script text; when given,
    each slide's speaker-notes pane is populated so reviewers see the
    script alongside the slide in PowerPoint's notes view, and downstream
    TTS can pull the text from the .pptx if it ever loses sync with
    script.md.

    `template_vars` (optional) maps NAME -> value; before writing any
    field, occurrences of "{{NAME}}" in field values (strings or list
    elements) are substituted. Used by `papercast approve` to fill
    `{{REPORT_DATE}}` on the Cover slide once the reviewer commits a
    date. Unset variables leave the "{{NAME}}" literal in place — that's
    the signal that the human gate hasn't been crossed yet.
    """
    template_path = Path(template_path)
    figures_dir = Path(figures_dir)
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(template_path)
    layouts_by_name = {layout.name: layout for layout in prs.slide_layouts}
    image_lookup = _load_figures_lookup(figures_dir)
    tvars = template_vars or {}

    for page in plan.pages:
        layout = layouts_by_name.get(page.layout)
        if layout is None:
            raise ValueError(
                f"layout {page.layout!r} not found in template; "
                f"available: {sorted(layouts_by_name)}"
            )
        slide = prs.slides.add_slide(layout)
        resolved = _resolve_page_fields(page, tvars)
        _fill_page(slide, resolved, layout, image_lookup, figures_dir)
        if page_notes:
            note_text = page_notes.get(page.page_no)
            if note_text:
                slide.notes_slide.notes_text_frame.text = note_text

    prs.save(out_path)


def _resolve_page_fields(page: PageSpec, tvars: dict[str, str]) -> PageSpec:
    """Substitute {{NAME}} placeholders inside page field values. Returns
    a new PageSpec — the input is left untouched so callers can re-use
    the same plan with different template_vars."""
    if not tvars:
        return page
    new_fields: dict[str, Any] = {}
    for k, v in page.fields.items():
        if isinstance(v, str):
            new_fields[k] = _substitute(v, tvars)
        elif isinstance(v, list):
            new_fields[k] = [
                _substitute(item, tvars) if isinstance(item, str) else item
                for item in v
            ]
        else:
            new_fields[k] = v
    return PageSpec(page_no=page.page_no, layout=page.layout, fields=new_fields)


_TEMPLATE_VAR_RE = re.compile(r"\{\{(\w+)\}\}")


def _substitute(text: str, tvars: dict[str, str]) -> str:
    def repl(m: re.Match[str]) -> str:
        name = m.group(1)
        return tvars.get(name, m.group(0))  # leave literal if not provided
    return _TEMPLATE_VAR_RE.sub(repl, text)


# Header line "## Page N" splits the script.md into per-page sections.
_SCRIPT_PAGE_HEADER_RE = re.compile(r"^##\s*Page\s+(\d+)\s*$", re.MULTILINE)

# A line of three or more dashes on its own marks the start of an optional
# metadata fence at the END of the document (e.g. `total_chars: 1872`).
# When present inside what would otherwise be the last page's notes, we
# strip it so the script doesn't leak metadata into the speaker-notes pane
# (which the TTS would also read aloud).
_METADATA_FENCE_RE = re.compile(r"^-{3,}\s*$", re.MULTILINE)


def parse_script_md(path: Path) -> dict[int, str]:
    """Read script.md and return {page_no -> spoken text}.

    Format expected (matches what the Author Agent / hand-author writes):

        # Optional title (ignored)
        > Optional metadata blockquote (ignored)

        ## Page 1
        ...spoken text for slide 1, possibly multi-line...

        ## Page 2
        ...

        ---
        total_chars: 1234              <- optional metadata fence; stripped
        estimated_seconds: 320

    Headers other than `## Page N` and any text before the first `## Page`
    header are dropped. A trailing `---` fence and everything after it is
    treated as metadata and stripped from the LAST page's notes so the
    .pptx speaker notes never contain `total_chars: ...` etc. Missing
    file returns an empty dict so the caller can treat notes as optional.
    """
    path = Path(path)
    if not path.exists():
        return {}
    content = path.read_text(encoding="utf-8")
    notes: dict[int, str] = {}
    matches = list(_SCRIPT_PAGE_HEADER_RE.finditer(content))
    for i, m in enumerate(matches):
        page_no = int(m.group(1))
        body_start = m.end()
        body_end = matches[i + 1].start() if i + 1 < len(matches) else len(content)
        body = content[body_start:body_end]
        # On the last page, strip an optional `---`-led metadata fence
        # so `total_chars: ...` lines don't bleed into speaker notes.
        if i == len(matches) - 1:
            fence = _METADATA_FENCE_RE.search(body)
            if fence:
                body = body[: fence.start()]
        body = body.strip()
        if body:
            notes[page_no] = body
    return notes


# ---------------------------------------------------------------------------
# Internals
# ---------------------------------------------------------------------------


def _load_figures_lookup(figures_dir: Path) -> dict[str, str]:
    """Return {figure_id -> filename}. If figures.json is missing we tolerate
    that — slides without Image fields will still render. assemble_pptx will
    raise later if a plan field references a missing id."""
    meta = figures_dir / "figures.json"
    if not meta.exists():
        return {}
    payload = json.loads(meta.read_text(encoding="utf-8"))
    return {f["id"]: f["filename"] for f in payload}


def _placeholders_by_name(slide_or_layout) -> dict[str, Any]:
    """Build a {name -> placeholder} lookup for a slide or layout."""
    return {ph.name: ph for ph in slide_or_layout.placeholders}


def _fill_page(
    slide,
    page: PageSpec,
    layout: SlideLayout,
    image_lookup: dict[str, str],
    figures_dir: Path,
) -> None:
    """Map plan field-names to placeholder idx via the layout (the layout
    is where semantic names live), then look up the slide placeholder by
    idx. python-pptx auto-renames slide-level placeholders on insertion,
    so name-based lookup on the slide itself is unreliable."""
    # name -> idx, derived from the layout (authoritative for naming).
    name_to_idx: dict[str, int] = {}
    for ph in layout.placeholders:
        name_to_idx[ph.name] = ph.placeholder_format.idx
    # idx -> placeholder, on the slide.
    slide_phs_by_idx: dict[int, Any] = {
        ph.placeholder_format.idx: ph for ph in slide.placeholders
    }

    for field_name, value in page.fields.items():
        idx = name_to_idx.get(field_name)
        if idx is None:
            # Layout does not define this field — skip silently so plans
            # remain portable across templates with extra optional fields.
            continue
        ph = slide_phs_by_idx.get(idx)
        if ph is None:
            continue  # cloned placeholder somehow missing from slide
        if field_name == "Image" or _looks_like_image_field(field_name, value):
            _fill_image(ph, value, image_lookup, figures_dir, slide)
            continue
        if isinstance(value, list):
            _fill_text_paragraphs(ph, [str(v) for v in value])
        else:
            _fill_text_single(ph, str(value))
        _apply_field_styling(ph, page.layout, field_name)


def _apply_field_styling(ph, layout_name: str, field_name: str) -> None:
    """Apply per-field visual conventions that the master doesn't enforce
    yet. These are intentionally narrow — extending them is cheaper than
    re-editing the lab template every iteration.
    """
    # Cover titles: center alignment.
    if layout_name == "Cover" and field_name in ("Title", "Title_chinese"):
        for p in ph.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER

    # Subtitle (small label, upper-left of every content layout): bold.
    if field_name == "Subtitle":
        for p in ph.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True

    # Bullets — adaptive font size + autosize fallback.
    # The master assumes ~5 short bullets; LLM plans frequently produce
    # 6-10 longer bullets, which then overflow into the image below.
    if field_name == "Bullets":
        _clamp_bullets_font_size(ph)


def _clamp_bullets_font_size(ph) -> None:
    """Pick a font size for the Bullets placeholder based on paragraph
    count, then enable PowerPoint's "shrink text on overflow" as a
    secondary safety net for unusually long lines.

    Sizing schedule (chosen against the lab template's 2.9-cm Bullets
    box; tested visually with the FPC-VLA paper):

        n ≤ 5    → 18 pt   (master default; no change)
        6–7      → 16 pt
        8–9      → 14 pt
        ≥ 10     → 12 pt

    We set the font size on every run (not just the first) because
    python-pptx creates one default run per paragraph and we want the
    schedule to apply uniformly.
    """
    tf = ph.text_frame
    n = max(1, len(tf.paragraphs))
    if n <= 5:
        size = 18
    elif n <= 7:
        size = 16
    elif n <= 9:
        size = 14
    else:
        size = 12

    pt_size = Pt(size)
    for para in tf.paragraphs:
        # If the paragraph has no run (text was set via tf.text or
        # paragraph.text), accessing .runs returns []; in that case the
        # paragraph carries default formatting and we set it via the
        # first run we can find — falling back to no-op silently if the
        # paragraph genuinely has no text content.
        if para.runs:
            for run in para.runs:
                run.font.size = pt_size
        elif para.text:
            # python-pptx exposes a default run via para.font; setting on
            # paragraph-level font cascades to its (implicit) runs.
            para.font.size = pt_size

    # Belt-and-braces: let PowerPoint shrink further if our schedule was
    # still too generous for the actual rendered width.
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.word_wrap = True
    except Exception:
        # Some placeholder kinds (rare) reject auto_size; ignore — the
        # explicit pt_size above is already in effect.
        pass


def _looks_like_image_field(name: str, value: Any) -> bool:
    """Some image placeholders may not literally be named 'Image' — e.g. a
    template could expose 'Diagram'. Treat any field whose value is a string
    matching a known figure id as an image."""
    return False  # current templates only use 'Image'; reserved for the future


def _fill_text_single(ph, text: str) -> None:
    ph.text_frame.text = text


def _fill_text_paragraphs(ph, lines: list[str]) -> None:
    """Set the placeholder's text frame to one paragraph per line.

    python-pptx's text_frame already has one default paragraph; we reuse it
    for the first line then append the rest.
    """
    tf = ph.text_frame
    tf.clear()
    if not lines:
        return
    tf.paragraphs[0].text = lines[0]
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.text = extra


def _fill_image(
    ph,
    value: Any,
    image_lookup: dict[str, str],
    figures_dir: Path,
    slide,
) -> None:
    """Insert an image into the placeholder using *contain* fit:

    - keep the image's original aspect ratio
    - scale it to fit inside the placeholder's bounding box
    - center it; let the unused dimension be empty space (letterbox /
      pillarbox), never crop and never stretch.

    python-pptx's `placeholder.insert_picture()` does crop-to-fill on
    PICTURE-type placeholders and stretch-to-fill on OBJECT placeholders,
    neither of which is acceptable for scientific figures. So we always
    delete the placeholder and add a free-floating Picture sized to the
    contain-fit geometry.
    """
    if not isinstance(value, str) or not value:
        raise ValueError(f"image field expects a figure id string, got {value!r}")
    filename = image_lookup.get(value)
    if filename is None:
        raise ValueError(
            f"image {value!r} not found in figures.json under {figures_dir}"
        )
    img_path = figures_dir / filename
    if not img_path.exists():
        raise ValueError(f"image file missing on disk: {img_path}")

    box_left, box_top = ph.left, ph.top
    box_w, box_h = ph.width, ph.height

    # Read image dimensions in pixels — only the ratio matters here.
    from PIL import Image as _PILImage
    with _PILImage.open(img_path) as img:
        img_w_px, img_h_px = img.size

    new_left, new_top, new_w, new_h = _contain_fit(
        box_left, box_top, box_w, box_h, img_w_px, img_h_px
    )

    # Drop the placeholder and replace with a free-floating picture so
    # we have full control over the geometry.
    sp = ph._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(
        str(img_path), new_left, new_top, width=new_w, height=new_h,
    )


def _contain_fit(
    box_left: int, box_top: int, box_w: int, box_h: int,
    img_w: int, img_h: int,
) -> tuple[int, int, int, int]:
    """Return (left, top, width, height) for an image of `img_w x img_h`
    fitted inside the box at `(box_left, box_top, box_w, box_h)` while
    preserving aspect ratio.

    Box dimensions are in EMU; image dimensions are in pixels — but only
    their ratios feed into the calculation, so units don't have to match.
    """
    if img_h <= 0 or box_h <= 0:
        return box_left, box_top, box_w, box_h
    box_aspect = box_w / box_h
    img_aspect = img_w / img_h
    if img_aspect > box_aspect:
        # Image is wider — fit to box width, letterbox top/bottom.
        new_w = box_w
        new_h = int(box_w / img_aspect)
        new_left = box_left
        new_top = box_top + (box_h - new_h) // 2
    else:
        # Image is taller — fit to box height, pillarbox left/right.
        new_h = box_h
        new_w = int(box_h * img_aspect)
        new_top = box_top
        new_left = box_left + (box_w - new_w) // 2
    return new_left, new_top, new_w, new_h
