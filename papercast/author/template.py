"""Parse the lab PPT template into a JSON schema for the Author Agent.

Design (from §7.0 of the design doc):

    template (.pptx) + demo (.pptx) ──► lab_template.meta.json

The Author Agent never opens the .pptx at generation time; it produces JSON
matching the schema_examples shape and python-pptx assembles by writing into
placeholders by `idx`. That decoupling is what this module exists to enforce.
"""

from __future__ import annotations

import hashlib
import json
from dataclasses import asdict, dataclass, field
from datetime import UTC, datetime
from enum import StrEnum
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.slide import SlideLayout

PARSER_VERSION = "1.0"

EMU_PER_CM = 360000.0

_THEME_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


class ImageFit(StrEnum):
    CONTAIN = "contain"
    COVER = "cover"


@dataclass(frozen=True)
class PlaceholderSpec:
    idx: int
    name: str
    type: str  # python-pptx PP_PLACEHOLDER name (e.g. BODY, PICTURE, OBJECT)
    left_cm: float
    top_cm: float
    width_cm: float
    height_cm: float
    aspect: float | None = None
    fit: ImageFit | None = None


@dataclass(frozen=True)
class LayoutSpec:
    name: str
    placeholders: list[PlaceholderSpec]


@dataclass(frozen=True)
class ThemeSpec:
    font_major_latin: str
    font_minor_latin: str
    font_major_ea: str
    font_minor_ea: str
    colors: dict[str, str]


@dataclass
class TemplateMeta:
    template_sha1: str
    parsed_at: str
    parser_version: str
    slide_size_cm: tuple[float, float]
    theme: ThemeSpec
    layouts: list[LayoutSpec]
    schema_examples: dict[str, dict[str, str]] = field(default_factory=dict)
    layouts_without_examples: list[str] = field(default_factory=list)

    def layout_by_name(self, name: str) -> LayoutSpec | None:
        for layout in self.layouts:
            if layout.name == name:
                return layout
        return None


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def parse_template(template_path: Path, demo_path: Path | None = None) -> TemplateMeta:
    """Read a .pptx template and produce TemplateMeta.

    `demo_path` (optional) is a sibling .pptx with one filled slide per layout.
    Its placeholder text becomes `schema_examples` so the LLM has a concrete
    example of what each layout's content looks like.
    """
    template_path = Path(template_path)
    if not template_path.exists():
        raise FileNotFoundError(f"template not found: {template_path}")

    sha1 = _file_sha1(template_path)
    prs = Presentation(template_path)

    slide_size_cm = (prs.slide_width / EMU_PER_CM, prs.slide_height / EMU_PER_CM)
    theme = _extract_theme(prs)
    layouts = [_layout_to_spec(layout) for layout in prs.slide_layouts]

    schema_examples: dict[str, dict[str, str]] = {}
    if demo_path is not None and Path(demo_path).exists():
        schema_examples = _extract_examples(Path(demo_path))

    layouts_without_examples = [
        layout.name for layout in layouts if layout.name not in schema_examples
    ]

    return TemplateMeta(
        template_sha1=sha1,
        parsed_at=datetime.now(UTC).isoformat(timespec="seconds"),
        parser_version=PARSER_VERSION,
        slide_size_cm=slide_size_cm,
        theme=theme,
        layouts=layouts,
        schema_examples=schema_examples,
        layouts_without_examples=layouts_without_examples,
    )


def write_meta(meta: TemplateMeta, out_path: Path) -> None:
    """Serialize TemplateMeta to JSON, in a stable, human-readable shape."""
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        json.dumps(_meta_to_json(meta), indent=2, ensure_ascii=False, sort_keys=False),
        encoding="utf-8",
    )


# ---------------------------------------------------------------------------
# Internals
# ---------------------------------------------------------------------------


def _file_sha1(path: Path) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 20), b""):
            h.update(chunk)
    return h.hexdigest()


def _placeholder_type_name(ph: Any) -> str:
    try:
        return str(ph.placeholder_format.type).split(".")[-1].split(" ")[0]
    except Exception:  # pragma: no cover - defensive
        return "UNKNOWN"


def _emu_to_cm(value: int | None) -> float:
    if value is None:
        return 0.0
    return round(value / EMU_PER_CM, 3)


def _layout_to_spec(layout: SlideLayout) -> LayoutSpec:
    placeholders: list[PlaceholderSpec] = []
    for ph in layout.placeholders:
        idx = ph.placeholder_format.idx
        name = ph.name
        ptype = _placeholder_type_name(ph)
        left = _emu_to_cm(ph.left)
        top = _emu_to_cm(ph.top)
        width = _emu_to_cm(ph.width)
        height = _emu_to_cm(ph.height)
        aspect = round(width / height, 3) if height > 0 else None

        is_image = name.lower().startswith("image") or ptype == "PICTURE"
        fit = ImageFit.CONTAIN if is_image else None

        placeholders.append(PlaceholderSpec(
            idx=idx,
            name=name,
            type=ptype,
            left_cm=left,
            top_cm=top,
            width_cm=width,
            height_cm=height,
            aspect=aspect if is_image else None,
            fit=fit,
        ))
    return LayoutSpec(name=layout.name, placeholders=placeholders)


def _extract_theme(prs: Presentation) -> ThemeSpec:
    master = prs.slide_masters[0]
    theme_part = None
    for rel in master.part.rels.values():
        if "theme" in rel.reltype:
            theme_part = rel.target_part
            break
    if theme_part is None:  # pragma: no cover - every Office file has a theme
        return ThemeSpec("", "", "", "", {})

    root = etree.fromstring(theme_part.blob)

    def _font(parent_xpath: str, child: str) -> str:
        node = root.find(f"{parent_xpath}/a:{child}", _THEME_NS)
        return node.get("typeface", "") if node is not None else ""

    colors: dict[str, str] = {}
    for color_node in root.findall(".//a:clrScheme/*", _THEME_NS):
        slot = etree.QName(color_node).localname
        for child in color_node:
            ctag = etree.QName(child).localname
            val: str | None = None
            if ctag == "srgbClr":
                val = "#" + child.get("val", "").upper()
            elif ctag == "sysClr":
                last = child.get("lastClr")
                if last:
                    val = "#" + last.upper()
            if val is not None:
                colors[slot] = val
                break

    return ThemeSpec(
        font_major_latin=_font(".//a:fontScheme/a:majorFont", "latin"),
        font_minor_latin=_font(".//a:fontScheme/a:minorFont", "latin"),
        font_major_ea=_font(".//a:fontScheme/a:majorFont", "ea"),
        font_minor_ea=_font(".//a:fontScheme/a:minorFont", "ea"),
        colors=colors,
    )


def _extract_examples(demo_path: Path) -> dict[str, dict[str, str]]:
    """Walk every demo slide; for each, capture {placeholder_name: text} keyed
    by layout name. If a layout appears multiple times, later slides only fill
    in fields the earlier ones missed."""
    prs = Presentation(demo_path)
    examples: dict[str, dict[str, str]] = {}

    for slide in prs.slides:
        layout_name = slide.slide_layout.name
        # Build idx -> name from the layout (slide placeholders may carry old
        # default names like "内容占位符 2" — the layout is the authoritative
        # source for naming).
        layout_idx_to_name = {
            ph.placeholder_format.idx: ph.name
            for ph in slide.slide_layout.placeholders
        }
        bucket = examples.setdefault(layout_name, {})
        for shp in slide.shapes:
            if not shp.is_placeholder:
                continue
            idx = shp.placeholder_format.idx
            ph_name = layout_idx_to_name.get(idx)
            if ph_name is None:
                continue  # orphan placeholder — skip
            if not shp.has_text_frame:
                continue
            text = shp.text_frame.text.strip()
            if not text:
                continue
            # First non-empty wins; later slides only fill in missing fields.
            bucket.setdefault(ph_name, text)

    return examples


def _meta_to_json(meta: TemplateMeta) -> dict[str, Any]:
    """Convert dataclasses to JSON-friendly dicts; preserve enum string values."""
    payload = asdict(meta)
    # `slide_size_cm` is a tuple — asdict turns it into a list, which is fine.
    # `fit` is an enum — asdict turns it into the enum object; coerce to string.
    for layout in payload["layouts"]:
        for ph in layout["placeholders"]:
            if ph.get("fit") is not None:
                ph["fit"] = ph["fit"].value if hasattr(ph["fit"], "value") else str(ph["fit"])
    return payload
